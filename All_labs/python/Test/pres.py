from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Создаем презентацию
prs = Presentation()

# Функция для стилизованного добавления текста
def add_styled_text(slide, text, font_size=32, bold=True, color=(0, 0, 0)):
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    p.alignment = PP_ALIGN.CENTER

# Слайд 1: Заголовок
slide_1 = prs.slides.add_slide(prs.slide_layouts[5])  # Заголовочный слайд
add_styled_text(slide_1, "Адамгершілік – асыл қасиет", font_size=44, color=(0, 102, 204))

# Слайд 2: Определение
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])  # Основной слайд
add_styled_text(slide_2, "Адамгершілік деген не?", font_size=36, color=(0, 0, 139))
slide_2.shapes.add_picture("path/to/image1.jpg", Inches(1), Inches(2), width=Inches(3), height=Inches(3))

# Слайд 3: Важность
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
add_styled_text(slide_3, "Адамгершілік қасиеттің маңызы", font_size=36, color=(0, 0, 139))
slide_3.shapes.add_picture("path/to/image2.jpg", Inches(5), Inches(2), width=Inches(3), height=Inches(3))

# Повторяйте добавление для других слайдов
# ...

# Сохранение презентации
prs.save("Адамгершілік_асыл_қасиет_стилизацияланған_презентация.pptx")
